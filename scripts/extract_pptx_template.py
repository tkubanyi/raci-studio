"""Extract presentation structure, theme, and shape map from a PPTX file."""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


def emu_to_inches(emu: int | None) -> float | None:
    if emu is None:
        return None
    return round(Emu(emu).inches, 4)


def emu_to_pt(emu: int | None) -> float | None:
    if emu is None:
        return None
    return round(Emu(emu).pt, 2)


def color_to_dict(color) -> dict | None:
    if color is None:
        return None
    try:
        if color.type is None:
            return None
    except Exception:
        return None

    result: dict = {"type": str(color.type)}

    try:
        if color.type == MSO_FILL.SOLID or hasattr(color, "rgb"):
            rgb = color.rgb
            if rgb is not None:
                result["rgb"] = f"#{rgb}"
                result["r"] = rgb[0]
                result["g"] = rgb[1]
                result["b"] = rgb[2]
    except Exception:
        pass

    try:
        if hasattr(color, "theme_color") and color.theme_color is not None:
            result["theme_color"] = str(color.theme_color)
    except Exception:
        pass

    try:
        if hasattr(color, "brightness") and color.brightness is not None:
            result["brightness"] = color.brightness
    except Exception:
        pass

    return result if len(result) > 1 else None


def fill_to_dict(fill) -> dict | None:
    if fill is None:
        return None
    try:
        fill_type = fill.type
    except Exception:
        return None

    result: dict = {"type": str(fill_type)}

    try:
        if fill_type == MSO_FILL.SOLID:
            result["color"] = color_to_dict(fill.fore_color)
    except Exception:
        pass

    return result if len(result) > 1 else {"type": str(fill_type)}


def line_to_dict(line) -> dict | None:
    if line is None:
        return None
    result: dict = {}
    try:
        if line.fill:
            result["fill"] = fill_to_dict(line.fill)
    except Exception:
        pass
    try:
        if line.width is not None:
            result["width_pt"] = emu_to_pt(line.width)
    except Exception:
        pass
    return result or None


def font_to_dict(font) -> dict:
    data: dict = {}
    if font.name:
        data["name"] = font.name
    if font.size is not None:
        data["size_pt"] = emu_to_pt(font.size)
    if font.bold is not None:
        data["bold"] = font.bold
    if font.italic is not None:
        data["italic"] = font.italic
    if font.underline is not None:
        data["underline"] = font.underline
    try:
        if font.color:
            data["color"] = color_to_dict(font.color)
    except Exception:
        pass
    return data


def paragraph_to_dict(paragraph) -> dict:
    data: dict = {"text": paragraph.text}
    if paragraph.level is not None:
        data["level"] = paragraph.level
    if paragraph.alignment is not None:
        data["alignment"] = str(paragraph.alignment)
    try:
        if paragraph.font:
            font_data = font_to_dict(paragraph.font)
            if font_data:
                data["font"] = font_data
    except Exception:
        pass
    runs = []
    for run in paragraph.runs:
        run_data: dict = {"text": run.text}
        try:
            if run.font:
                font_data = font_to_dict(run.font)
                if font_data:
                    run_data["font"] = font_data
        except Exception:
            pass
        if run_data.get("text"):
            runs.append(run_data)
    if runs:
        data["runs"] = runs
    return data


def text_frame_to_dict(text_frame) -> dict | None:
    if text_frame is None:
        return None
    paragraphs = [paragraph_to_dict(p) for p in text_frame.paragraphs if p.text.strip()]
    if not paragraphs:
        return None
    data: dict = {"paragraphs": paragraphs}
    try:
        if text_frame.margin_left is not None:
            data["margin_left_in"] = emu_to_inches(text_frame.margin_left)
        if text_frame.margin_right is not None:
            data["margin_right_in"] = emu_to_inches(text_frame.margin_right)
        if text_frame.margin_top is not None:
            data["margin_top_in"] = emu_to_inches(text_frame.margin_top)
        if text_frame.margin_bottom is not None:
            data["margin_bottom_in"] = emu_to_inches(text_frame.margin_bottom)
    except Exception:
        pass
    return data


def shape_position(shape) -> dict:
    return {
        "left_in": emu_to_inches(shape.left),
        "top_in": emu_to_inches(shape.top),
        "width_in": emu_to_inches(shape.width),
        "height_in": emu_to_inches(shape.height),
    }


def shape_to_dict(shape, depth: int = 0) -> dict:
    data: dict = {
        "id": shape.shape_id,
        "name": shape.name,
        "type": str(shape.shape_type),
        "position": shape_position(shape),
    }

    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        try:
            data["auto_shape_type"] = str(shape.auto_shape_type)
        except Exception:
            pass

    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        try:
            data["placeholder_format"] = {
                "idx": shape.placeholder_format.idx,
                "type": str(shape.placeholder_format.type),
            }
        except Exception:
            pass

    try:
        if shape.fill:
            fill_data = fill_to_dict(shape.fill)
            if fill_data:
                data["fill"] = fill_data
    except Exception:
        pass

    try:
        if shape.line:
            line_data = line_to_dict(shape.line)
            if line_data:
                data["line"] = line_data
    except Exception:
        pass

    try:
        if shape.has_text_frame:
            text_data = text_frame_to_dict(shape.text_frame)
            if text_data:
                data["text"] = text_data
    except Exception:
        pass

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        children = []
        for child in shape.shapes:
            children.append(shape_to_dict(child, depth + 1))
        data["children"] = children

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            data["image"] = {
                "content_type": shape.image.content_type,
                "filename": shape.image.filename,
            }
        except Exception:
            pass

    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table = shape.table
        data["table"] = {
            "rows": len(table.rows),
            "cols": len(table.columns),
            "cells": [],
        }
        for r_idx, row in enumerate(table.rows):
            for c_idx, cell in enumerate(row.cells):
                cell_data = {
                    "row": r_idx,
                    "col": c_idx,
                    "text": cell.text.strip(),
                }
                try:
                    if cell.fill:
                        cell_fill = fill_to_dict(cell.fill)
                        if cell_fill:
                            cell_data["fill"] = cell_fill
                except Exception:
                    pass
                data["table"]["cells"].append(cell_data)

    return data


def extract_theme_colors(pptx_path: Path) -> dict:
    import xml.etree.ElementTree as ET
    import zipfile

    colors: dict = {}
    try:
        with zipfile.ZipFile(pptx_path) as z:
            theme_files = sorted(
                n for n in z.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml")
            )
            if not theme_files:
                return colors
            root = ET.fromstring(z.read(theme_files[0]))
            ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            clr_scheme = root.find(".//a:clrScheme", ns)
            if clr_scheme is not None:
                for child in clr_scheme:
                    tag = child.tag.split("}")[-1]
                    srgb = child.find("a:srgbClr", ns)
                    sys_clr = child.find("a:sysClr", ns)
                    if srgb is not None and srgb.get("val"):
                        colors[tag] = f"#{srgb.get('val')}"
                    elif sys_clr is not None and sys_clr.get("lastClr"):
                        colors[tag] = f"#{sys_clr.get('lastClr')}"
    except Exception as exc:
        colors["_error"] = str(exc)
    return colors


def extract_slide_layouts(prs: Presentation) -> list[dict]:
    layouts = []
    seen = set()
    for slide in prs.slides:
        layout = slide.slide_layout
        key = layout.name
        if key in seen:
            continue
        seen.add(key)
        layouts.append(
            {
                "name": layout.name,
                "placeholders": [
                    {
                        "idx": ph.placeholder_format.idx,
                        "type": str(ph.placeholder_format.type),
                        "name": ph.name,
                    }
                    for ph in layout.placeholders
                ],
            }
        )
    return layouts


def summarize_colors(shapes: list[dict]) -> dict:
    fills = set()
    lines = set()
    fonts = set()

    def walk(shape_list):
        for shape in shape_list:
            fill = shape.get("fill", {})
            color = fill.get("color", {}) if fill else {}
            if color.get("rgb"):
                fills.add(color["rgb"])
            if color.get("theme_color"):
                fills.add(f"theme:{color['theme_color']}")

            line = shape.get("line", {})
            line_fill = line.get("fill", {}) if line else {}
            line_color = line_fill.get("color", {}) if line_fill else {}
            if line_color.get("rgb"):
                lines.add(line_color["rgb"])

            text = shape.get("text", {})
            for para in text.get("paragraphs", []):
                font = para.get("font", {})
                fc = font.get("color", {})
                if fc and fc.get("rgb"):
                    fonts.add(fc["rgb"])
                for run in para.get("runs", []):
                    rfc = run.get("font", {}).get("color", {})
                    if rfc and rfc.get("rgb"):
                        fonts.add(rfc["rgb"])

            walk(shape.get("children", []))

    walk(shapes)
    return {
        "fill_colors": sorted(fills),
        "line_colors": sorted(lines),
        "font_colors": sorted(fonts),
    }


def infer_slide_role(slide_data: dict) -> str:
    texts = []
    for shape in slide_data.get("shapes", []):
        text = shape.get("text", {})
        for para in text.get("paragraphs", []):
            if para.get("text"):
                texts.append(para["text"].strip())

    joined = " ".join(texts).lower()
    if slide_data["index"] == 1:
        return "title_cover"
    if "agenda" in joined or "contents" in joined:
        return "agenda"
    if "thank" in joined or "questions" in joined or "contact" in joined:
        return "closing"
    if any(shape.get("type") == "TABLE (19)" for shape in slide_data.get("shapes", [])):
        return "table"
    if len(slide_data.get("shapes", [])) >= 4:
        return "content_multi_block"
    return "content"


def build_template_map(pptx_path: Path) -> dict:
    prs = Presentation(str(pptx_path))

    slide_width_in = emu_to_inches(prs.slide_width)
    slide_height_in = emu_to_inches(prs.slide_height)

    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        shapes = [shape_to_dict(shape) for shape in slide.shapes]
        slide_data = {
            "index": idx,
            "layout_name": slide.slide_layout.name,
            "shapes": shapes,
            "shape_count": len(shapes),
            "color_summary": summarize_colors(shapes),
        }
        slide_data["role"] = infer_slide_role(slide_data)
        slides.append(slide_data)

    all_fills = set()
    all_lines = set()
    all_fonts = set()
    for slide in slides:
        for key, target in [
            ("fill_colors", all_fills),
            ("line_colors", all_lines),
            ("font_colors", all_fonts),
        ]:
            target.update(slide["color_summary"][key])

    return {
        "source_file": pptx_path.name,
        "source_path": str(pptx_path),
        "version": "1.0",
        "purpose": "Template map for replicating presentation format with new narratives",
        "canvas": {
            "width_in": slide_width_in,
            "height_in": slide_height_in,
            "aspect_ratio": "16:9" if slide_width_in and slide_height_in and slide_width_in > slide_height_in else "4:3",
        },
        "theme": {
            "colors": extract_theme_colors(pptx_path),
        },
        "layouts_used": extract_slide_layouts(prs),
        "palette": {
            "fill_colors": sorted(all_fills),
            "line_colors": sorted(all_lines),
            "font_colors": sorted(all_fonts),
        },
        "slide_count": len(slides),
        "slides": slides,
        "replication_notes": {
            "instructions": [
                "Use canvas dimensions and theme colors as fixed design tokens.",
                "Match each new slide to the closest role (title_cover, agenda, content, table, closing).",
                "Preserve shape positions, auto_shape_type, fill, and line styling from the mapped slide.",
                "Replace text content in paragraphs/runs while keeping font sizes and colors.",
                "Grouped shapes should be recreated as groups with the same child structure.",
            ],
            "roles": {
                "title_cover": "Opening slide with title and subtitle branding",
                "agenda": "Table of contents / section overview",
                "content": "Standard narrative slide with text blocks",
                "content_multi_block": "Slide with multiple content zones or icons",
                "table": "Tabular data layout",
                "closing": "Thank you / Q&A / contact slide",
            },
        },
    }


def main() -> int:
    if len(sys.argv) < 3:
        print("Usage: python extract_pptx_template.py <input.pptx> <output.json>")
        return 1

    input_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])

    if not input_path.exists():
        print(f"File not found: {input_path}")
        return 1

    template_map = build_template_map(input_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(template_map, indent=2, ensure_ascii=False), encoding="utf-8")

    print(f"Extracted {template_map['slide_count']} slides")
    print(f"Theme colors: {len(template_map['theme']['colors'])}")
    print(f"Saved to: {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
