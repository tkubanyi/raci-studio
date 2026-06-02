"""Audit text placement in a generated presentation."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def walk(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk(shape.shapes)
        else:
            yield shape


def main() -> int:
    path = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(
        r"c:\Users\tkubanyi001\Projects\raci-studio\data\presentation-templates\Project_Vienna_Discovery_Phase.pptx"
    )
    prs = Presentation(str(path))
    for si, slide in enumerate(prs.slides, 1):
        print(f"--- Slide {si} ---")
        for s in walk(slide.shapes):
            if not s.has_text_frame:
                continue
            t = s.text.strip()
            if not t:
                continue
            w, h = round(s.width / 914400, 2), round(s.height / 914400, 2)
            fs = set()
            for para in s.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        fs.add(int(run.font.size.pt))
            preview = t[:70].replace("\n", " | ")
            print(f"  {s.name[:32]:32} {w:5.2f}x{h:5.2f} fs={sorted(fs) or '?'} len={len(t):3d} | {preview}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
