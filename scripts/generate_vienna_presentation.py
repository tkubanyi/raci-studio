"""Generate Project Vienna presentation with template-accurate text fitting."""

from __future__ import annotations

import json
import re
import shutil
from datetime import date
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent.parent
TEMPLATE_PPTX = BASE / "data" / "presentation-templates" / "Proposal_Global_Payments_V2.pptx"
SUMMARY_JSON = BASE / "data" / "presentation-templates" / "proposal_global_payments_v3.summary.json"
DOCX_DIR = Path(r"C:\Users\tkubanyi001\OneDrive - PwC\Documents\Customers\Global Payments")
DOCX_PATH = DOCX_DIR / "Project Vienna_discovery phase.docx"
OUTPUT_PPTX = DOCX_DIR / "Project_Vienna_Discovery_Phase.pptx"
OUTPUT_FALLBACK = BASE / "data" / "presentation-templates" / "Project_Vienna_Discovery_Phase.pptx"

FOOTER = "Project Vienna — Diagnostic and Discovery Phase | Global Payments Europe"

# Template typography from proposal_global_payments_v3.summary.json
FONT_TITLE = 16.0
FONT_SECTION = 15.0
FONT_BODY = 12.0
FONT_SMALL = 10.5
FONT_STAT_TITLE = 20.0
FONT_WEEK = 20.0
FONT_COVER = 28.0


def load_summary() -> dict:
    return json.loads(SUMMARY_JSON.read_text(encoding="utf-8"))


def parse_docx(path: Path) -> dict:
    doc = Document(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    title = paragraphs[0] if paragraphs else "Project Vienna"
    subtitle = paragraphs[1] if len(paragraphs) > 1 else ""

    sections: list[dict] = []
    current: dict | None = None

    for text in paragraphs[2:]:
        if re.match(r"^\d+\.\s", text) and len(text) < 120:
            if current:
                sections.append(current)
            current = {"heading": text, "paragraphs": [], "bullets": []}
            continue
        if current is None:
            continue
        if re.match(
            r"^[A-Z][A-Za-z /&]+(?:Mapping|Review|Logic|Inputs|Team|Approach|Assessment|Techniques)$",
            text,
        ):
            current["bullets"].append({"title": text, "body": ""})
        elif current["bullets"] and not current["bullets"][-1]["body"]:
            current["bullets"][-1]["body"] = text
        else:
            current["paragraphs"].append(text)

    if current:
        sections.append(current)

    timeline = []
    if doc.tables:
        for row in doc.tables[0].rows[1:]:
            timeline.append([c.text.strip() for c in row.cells])

    return {"title": title, "subtitle": subtitle, "sections": sections, "timeline": timeline}


def inches(shape) -> tuple[float, float, float, float]:
    return (
        round(shape.left / 914400, 3),
        round(shape.top / 914400, 3),
        round(shape.width / 914400, 3),
        round(shape.height / 914400, 3),
    )


def iter_shapes(shapes, include_groups: bool = True):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP and include_groups:
            yield from iter_shapes(shape.shapes)
        yield shape


def find_shape(slide, name: str):
    for shape in iter_shapes(slide.shapes):
        if shape.name == name:
            return shape
    return None


def find_shapes(slide, predicate) -> list:
    found = [s for s in iter_shapes(slide.shapes) if predicate(s)]
    return sorted(found, key=lambda s: (s.top, s.left))


def find_title(slide):
    for shape in iter_shapes(slide.shapes):
        if shape.is_placeholder and "TITLE" in str(shape.placeholder_format.type):
            return shape
    return find_shape(slide, "Title 1")


def _sample_font(shape, default: float = FONT_BODY) -> dict:
    if not shape.has_text_frame:
        return {"name": "Arial", "size_pt": default, "bold": False, "color": None}
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            font = run.font
            color = None
            try:
                if font.color and font.color.rgb:
                    color = font.color.rgb
            except Exception:
                pass
            return {
                "name": font.name or "Arial",
                "size_pt": font.size.pt if font.size else default,
                "bold": bool(font.bold),
                "color": color,
            }
    return {"name": "Arial", "size_pt": default, "bold": False, "color": None}


def _apply_font(run, font_spec: dict, size_pt: float | None = None) -> None:
    run.font.name = font_spec.get("name") or "Arial"
    run.font.size = Pt(size_pt or font_spec.get("size_pt") or FONT_BODY)
    run.font.bold = font_spec.get("bold", False)
    color = font_spec.get("color")
    if color:
        try:
            run.font.color.rgb = color
        except Exception:
            pass


def _clear_text_frame(tf) -> None:
    tf.clear()


def _first_paragraph(tf):
    if tf.paragraphs:
        return tf.paragraphs[0]
    return tf.add_paragraph()


def set_single_text(
    shape,
    text: str,
    *,
    font_pt: float | None = None,
    auto_fit: bool = True,
) -> None:
    if not shape or not shape.has_text_frame:
        return
    font_spec = _sample_font(shape, font_pt or FONT_BODY)
    configure_text_frame(shape, auto_fit=auto_fit)
    tf = shape.text_frame
    _clear_text_frame(tf)
    p = _first_paragraph(tf)
    p.text = text
    if p.runs:
        _apply_font(p.runs[0], font_spec, font_pt or font_spec["size_pt"])


def set_multiline_text(
    shape,
    lines: list[str],
    *,
    font_pt: float = FONT_BODY,
    header_pt: float | None = None,
    bullet_from: int = 0,
    auto_fit: bool = True,
) -> None:
    if not shape or not shape.has_text_frame:
        return
    font_spec = _sample_font(shape, font_pt)
    configure_text_frame(shape, auto_fit=auto_fit)
    tf = shape.text_frame
    _clear_text_frame(tf)

    first = True
    for i, line in enumerate(lines):
        if not line:
            continue
        p = _first_paragraph(tf) if first else tf.add_paragraph()
        first = False
        is_bullet = bullet_from >= 0 and i >= bullet_from
        p.level = 1 if is_bullet else 0
        p.text = f"• {line}" if is_bullet else line
        size = header_pt if i == 0 and header_pt else font_pt
        if p.runs:
            _apply_font(p.runs[0], font_spec, size)
            if i == 0 and header_pt:
                p.runs[0].font.bold = True


def configure_text_frame(shape, *, anchor=MSO_ANCHOR.TOP, auto_fit: bool = True) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if auto_fit else MSO_AUTO_SIZE.NONE
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass


def truncate(text: str, max_chars: int) -> str:
    text = re.sub(r"\s+", " ", text.strip())
    if len(text) <= max_chars:
        return text
    cut = text[: max_chars - 1].rsplit(" ", 1)[0]
    return cut + "…"


def force_set_text(shape, text: str, *, font_pt: float = FONT_BODY, auto_fit: bool = True) -> None:
    """Replace all text in a shape — reliable for placeholders."""
    if not shape or not shape.has_text_frame:
        return
    font_spec = _sample_font(shape, font_pt)
    configure_text_frame(shape, auto_fit=auto_fit)
    shape.text = text
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            _apply_font(run, font_spec, font_pt)


def force_set_multiline(
    shape,
    lines: list[str],
    *,
    font_pt: float = FONT_BODY,
    header_pt: float | None = None,
    bullet_from: int = -1,
) -> None:
    if not shape or not shape.has_text_frame:
        return
    configure_text_frame(shape, auto_fit=True)
    shape.text = "\n".join(lines)
    font_spec = _sample_font(shape, font_pt)
    for i, para in enumerate(shape.text_frame.paragraphs):
        is_bullet = bullet_from >= 0 and i >= bullet_from
        if is_bullet and para.text and not para.text.startswith("•"):
            para.text = f"• {para.text}"
            para.level = 1
        size = header_pt if i == 0 and header_pt else font_pt
        for run in para.runs:
            _apply_font(run, font_spec, size)
            if i == 0 and header_pt:
                run.font.bold = True


def set_title_text(shape, text: str) -> None:
    if not shape:
        return
    configure_text_frame(shape, auto_fit=True)
    font_spec = _sample_font(shape, FONT_TITLE)
    tf = shape.text_frame
    _clear_text_frame(tf)
    parts = text.split("\n", 1)
    for i, part in enumerate(parts):
        p = _first_paragraph(tf) if i == 0 else tf.add_paragraph()
        p.text = part
        if p.runs:
            _apply_font(p.runs[0], font_spec, FONT_TITLE if i == 0 else FONT_SECTION)


def fill_table_cell(cell, text: str, *, font_pt: float = FONT_BODY, bold: bool = False) -> None:
    text = truncate(text, 180)
    tf = cell.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    _clear_text_frame(tf)
    p = _first_paragraph(tf)
    p.text = text
    if p.runs:
        p.runs[0].font.name = "Arial"
        p.runs[0].font.size = Pt(font_pt)
        p.runs[0].font.bold = bold


def fill_table(table, rows: list[list[str]], *, header: list[str] | None = None) -> None:
    if header and len(table.rows) > 0:
        for c_idx, value in enumerate(header):
            if c_idx < len(table.columns):
                fill_table_cell(table.cell(0, c_idx), value, font_pt=FONT_SECTION, bold=True)
    start = 1 if header else 0
    for r_idx, row_data in enumerate(rows):
        t_row = start + r_idx
        if t_row >= len(table.rows):
            break
        for c_idx in range(len(table.columns)):
            value = row_data[c_idx] if c_idx < len(row_data) else ""
            fill_table_cell(table.cell(t_row, c_idx), value, bold=(c_idx == 0))


def replace_footer(slide) -> None:
    for shape in iter_shapes(slide.shapes):
        if shape.is_placeholder and "FOOTER" in str(shape.placeholder_format.type):
            set_single_text(shape, FOOTER, font_pt=10.0, auto_fit=True)


def replace_slide_number(slide, number: int) -> None:
    for shape in iter_shapes(slide.shapes):
        if shape.is_placeholder and "SLIDE_NUMBER" in str(shape.placeholder_format.type):
            set_single_text(shape, str(number), font_pt=10.0, auto_fit=False)


def build_content(doc: dict) -> dict:
    what = next((s for s in doc["sections"] if "What" in s["heading"]), None)
    so_what = next((s for s in doc["sections"] if "So What" in s["heading"]), None)

    opportunities = []
    if so_what:
        paras = so_what["paragraphs"]
        i = 0
        while i < len(paras):
            p = paras[i]
            if "focus include" in p.lower() or "key areas" in p.lower():
                i += 1
                continue
            if i + 1 < len(paras) and len(p) < 80 and len(paras[i + 1]) > 40:
                opportunities.append([p, truncate(paras[i + 1], 180)])
                i += 2
                continue
            i += 1
        if len(opportunities) < 7:
            defaults = [
                ["Service & Process Transparency", "Complete inventory across SSC, JVs and HQ"],
                ["Reduce Duplication", "Identify parallel execution across entities"],
                ["JV Best Practices", "Assess scalable mature capabilities in JVs"],
                ["Service Placement Decisions", "Criteria for local vs. central vs. scale-up"],
                ["RACI & Governance", "Clarify accountability between JVs, SSC and HQ"],
                ["P&L & Cost Visibility", "Baseline cost, headcount and allocations"],
                ["Scalable JV Launch Model", "Inputs for operational playbook and template"],
            ]
            opportunities = defaults

    workstreams = [
        ("Service & Process Mapping", "Map services/processes across SSC, JVs and HQ."),
        ("RACI & Governance Mapping", "Define roles; resolve ownership gaps and decision rights."),
        ("P&L, Cost & Headcount Review", "Establish financial and workforce baseline."),
        ("Opportunity Assessment", "Evaluate placement: local, SSC, JV scale-up or central build."),
        ("Future-State Blueprint Inputs", "Playbook, deployment template and transition roadmap."),
    ]

    team_roles = [
        ("Executive Sponsor", "Strategic direction and priority confirmation"),
        ("Project Lead / PM", "Delivery plan, governance and reporting"),
        ("Process Excellence Lead", "Process mapping and opportunity prioritization"),
        ("Finance / FP&A Lead", "P&L, cost allocation and business case"),
        ("SSC Operations Lead", "Prague SSC capabilities, capacity and constraints"),
    ]

    outputs_left = [
        "Service and process inventory (SSC, JVs, HQ)",
        "Process maps for priority services",
        "RACI matrix and decision rights",
        "P&L, cost and headcount baseline",
    ]
    outputs_right = [
        "Duplication and fragmentation analysis",
        "JV best-practice scale-up assessment",
        "Service placement recommendations",
        "Playbook and deployment template inputs",
    ]

    methodologies = [
        ("Service Catalogue", "APQC-aligned taxonomy for like-for-like comparison."),
        ("SIPOC & Process Maps", "End-to-end boundaries, handoffs and controls."),
        ("Lean Six Sigma", "Waste, rework and variation diagnostics."),
        ("RACI / RAPID", "Workshop-based role and decision-rights clarity."),
        ("GBS Maturity Lens", "Standardization, automation, KPIs and scalability."),
    ]

    future_states = [
        ("A. Remain Local", "Regulatory, language or proximity requires JV delivery."),
        ("B. Transition to SSC", "Standardizable processes for Prague SSC."),
        ("C. Scale JV Capability", "Reuse mature JV service across entities."),
        ("D. Build Central Capability", "New scalable solution required."),
    ]

    raci_blocks = [
        ("Responsible", "Who performs the activity"),
        ("Accountable", "Who owns the outcome and decision"),
        ("Consulted", "Who provides input before action"),
        ("Informed", "Who is updated after the decision"),
    ]

    problem_main = truncate(
        what["paragraphs"][0] if what and what["paragraphs"] else "",
        200,
    )
    problem_role = truncate(
        (what["paragraphs"][2] if what and len(what.get("paragraphs", [])) > 2 else what["paragraphs"][-1])
        if what and what["paragraphs"]
        else "",
        180,
    )

    timeline_phases = doc.get("timeline") or [
        ["Mobilization", "Weeks 1–2", "Scope, governance, templates, stakeholder map"],
        ["Discovery", "Weeks 3–6", "Interviews, service/process mapping, systems inventory"],
        ["Analysis", "Weeks 7–10", "Duplication, best practices, cost drivers"],
        ["Future-State", "Weeks 11–14", "Prioritize transitions and placement logic"],
        ["Handover", "Weeks 15–16", "Findings, backlog, roadmap, playbook inputs"],
    ]

    return {
        1: {"title": f"{doc['title']}\n{doc['subtitle']}", "subtitle": f"Global Payments Europe\nPwC\n{date.today().strftime('%d. %B %Y')}"},
        3: {"title": "What? — Problem Statement\nEstablishing a Fact-Based Operating Model View", "body1": problem_main, "body2": problem_role},
        5: {"title": "So What? — Opportunities During Discovery\nKey Focus Areas", "table_rows": opportunities[:7]},
        12: {"title": "Now What?", "subtitle": "Discovery Approach, Team and Timeline"},
        16: {"title": "Discovery Approach\nFive Integrated Workstreams", "workstreams": workstreams},
        9: {"title": 'Tentative Timeline\n8–16 Week Diagnostic Program', "phases": timeline_phases[:5]},
        10: {"title": "Required Project Team", "roles": team_roles},
        6: {"title": "Expected Outputs\nEnd-of-Discovery Deliverables", "left": outputs_left, "right": outputs_right},
        7: {"title": "Industry Standards & Methodologies\nStructured Discovery Techniques", "blocks": methodologies},
        8: {"title": "GBS Maturity & Benchmarking\nComparison to Leading Practice", "subtitle": "Standardization • Centralization • Automation • KPIs • Controls • Scalability"},
        22: {"title": "Service Placement Options\nFour Future-State Paths", "boxes": future_states},
        17: {"title": "RACI & Operating Model Design\nRoles Across JVs, SSC and HQ", "blocks": raci_blocks},
        18: {"title": "SIPOC & Process Mapping\nFrom Documentation to Simplification", "intro": "Addressing complexity across JVs, SSC and HQ", "body": truncate("SIPOC maps Suppliers, Inputs, Process, Outputs and Customers. Detailed flowcharts capture handoffs, controls, rework and system touchpoints — revealing duplication before transition.", 420), "steps": ["Define process boundaries (SIPOC)", "Map handoffs and decision points", "Identify duplication and gaps", "Simplify before transition", "Validate with process owners"]},
        15: {"title": "Thank you", "subtitle": "© 2026 PricewaterhouseCoopers. All rights reserved.\nStrictly confidential — non-binding discovery overview."},
    }


def apply_slide_1(slide, content: dict) -> None:
    set_title_text(find_title(slide), content["title"])
    sub = find_shape(slide, "Subtitle 8")
    if sub:
        set_multiline_text(sub, content["subtitle"].split("\n"), font_pt=FONT_SECTION, bullet_from=-1, auto_fit=True)


def apply_slide_3(slide, content: dict) -> None:
    set_title_text(find_title(slide), content["title"])
    box5 = find_shape(slide, "TextBox 5")
    box7 = find_shape(slide, "TextBox 7")
    if box5:
        set_single_text(box5, content["body1"], font_pt=FONT_BODY, auto_fit=True)
    if box7:
        set_single_text(box7, content["body2"], font_pt=FONT_BODY, auto_fit=True)


def apply_slide_5(slide, content: dict) -> None:
    set_title_text(find_title(slide), content["title"])
    for shape in iter_shapes(slide.shapes):
        if shape.has_table:
            rows = [[r[0], r[1], ""] for r in content["table_rows"]]
            fill_table(
                shape.table,
                rows,
                header=["Opportunity Area", "Discovery Focus", ""],
            )
            break


def apply_slide_6(slide, content: dict) -> None:
    set_title_text(find_title(slide), content["title"])
    left = find_shape(slide, "TextBox 6")
    right = find_shape(slide, "TextBox 33")
    if left:
        set_multiline_text(left, ["Key Deliverables", *content["left"]], header_pt=FONT_SECTION, font_pt=FONT_BODY, bullet_from=1)
    if right:
        set_multiline_text(right, ["Analysis & Recommendations", *content["right"]], header_pt=FONT_SECTION, font_pt=FONT_BODY, bullet_from=1)
    for name in ("Rectangle 7", "Rectangle 32", "Rectangle 34", "Rectangle 35"):
        shape = find_shape(slide, name)
        if shape:
            set_single_text(shape, "", auto_fit=True)


def apply_slide_7(slide, content: dict) -> None:
    set_title_text(find_title(slide), content["title"])
    columns = find_shapes(
        slide,
        lambda s: s.has_text_frame and s.name.startswith("TextBox") and inches(s)[2] < 3 and inches(s)[3] > 3,
    )
    columns = sorted(columns, key=lambda s: s.left)[:5]
    for col, (title, body) in zip(columns, content["blocks"]):
        set_multiline_text(col, [title, body], header_pt=FONT_SECTION, font_pt=FONT_BODY, bullet_from=-1)


def apply_slide_8(slide, content: dict) -> None:
    set_title_text(find_title(slide), content["title"])
    box = find_shape(slide, "TextBox 9")
    if box:
        set_single_text(box, content["subtitle"], font_pt=FONT_SECTION, auto_fit=True)


def apply_slide_9(slide, content: dict) -> None:
    set_title_text(find_title(slide), content["title"])
    pentagons = sorted(
        [s for s in iter_shapes(slide.shapes) if s.has_text_frame and "Pentagon" in s.name],
        key=lambda s: s.left,
    )
    textboxes = sorted(
        [s for s in iter_shapes(slide.shapes) if s.has_text_frame and s.name.startswith("TextBox")],
        key=lambda s: (s.top, s.left),
    )
    phases = content["phases"]
    for i, phase in enumerate(phases):
        week_label = phase[1].replace("Weeks ", "Wk ").replace("–", "-")
        if i < len(pentagons):
            set_single_text(pentagons[i], week_label, font_pt=FONT_WEEK, auto_fit=True)
        if i < len(textboxes):
            set_multiline_text(
                textboxes[i],
                [phase[0], truncate(phase[2], 95)],
                header_pt=FONT_SECTION,
                font_pt=FONT_BODY,
                bullet_from=-1,
            )


def apply_slide_10(slide, content: dict) -> None:
    title = find_title(slide)
    if title:
        set_single_text(title, content["title"], font_pt=FONT_TITLE, auto_fit=True)
    cards = [
        find_shape(slide, "Google Shape;490;p66"),
        find_shape(slide, "Google Shape;484;p66"),
        find_shape(slide, "Google Shape;512;p66"),
        find_shape(slide, "Google Shape;485;p66"),
        find_shape(slide, "Google Shape;486;p66"),
    ]
    for card, (role, desc) in zip(cards, content["roles"]):
        if card:
            set_multiline_text(card, [role, desc], header_pt=FONT_SECTION, font_pt=FONT_BODY, bullet_from=-1)
    for name in ("Google Shape;487;p66",):
        shape = find_shape(slide, name)
        if shape:
            set_single_text(shape, "", auto_fit=True)


def apply_slide_12(slide, content: dict) -> None:
    title = find_title(slide)
    if title:
        set_title_text(title, content["title"])
    for shape in iter_shapes(slide.shapes):
        if shape.is_placeholder and "BODY" in str(shape.placeholder_format.type):
            set_single_text(shape, content["subtitle"], font_pt=FONT_COVER, auto_fit=True)
            return
    # Section header layout may only have title — append subtitle below title area
    if title:
        set_title_text(title, f"{content['title']}\n{content['subtitle']}")


def apply_slide_16(slide, content: dict) -> None:
    set_title_text(find_title(slide), content["title"])
    ws_rectangles = {
        "Rectangle 10": 0,
        "Rectangle 18": 1,
        "Rectangle 19": 2,
        "Rectangle 21": 3,
        "Rectangle 20": 4,
    }
    for rect_name, idx in ws_rectangles.items():
        shape = find_shape(slide, rect_name)
        if shape and idx < len(content["workstreams"]):
            title, desc = content["workstreams"][idx]
            set_multiline_text(
                shape,
                [f"WS{idx + 1}: {truncate(title, 40)}", truncate(desc, 70)],
                header_pt=10.0,
                font_pt=FONT_SMALL,
                bullet_from=-1,
            )
    # Update WS labels on diagram
    ws_labels = sorted(
        [
            s
            for s in iter_shapes(slide.shapes)
            if s.has_text_frame and s.name.startswith("Content Placeholder")
        ],
        key=lambda s: s.left,
    )
    for idx, label_shape in enumerate(ws_labels):
        if idx < 5:
            force_set_text(label_shape, f"WS{idx + 1}", font_pt=11.0)
        else:
            force_set_text(label_shape, "", font_pt=11.0)
    # Clear legacy template copy
    for name in ("Rectangle 22", "TextBox 38", "TextBox 30", "TextBox 40", "TextBox 41", "TextBox 42", "TextBox 17"):
        shape = find_shape(slide, name)
        if shape:
            set_single_text(shape, "", auto_fit=True)
    set_single_text(find_shape(slide, "TextBox 31"), "Discovery Phase", font_pt=18.0, auto_fit=True)


def apply_slide_17(slide, content: dict) -> None:
    set_title_text(find_title(slide), content["title"])
    placeholders = sorted(
        [
            s
            for s in iter_shapes(slide.shapes)
            if s.has_text_frame and "Content Placeholder" in s.name
        ],
        key=lambda s: s.left,
    )
    for ph, (label, desc) in zip(placeholders[:4], content["blocks"]):
        force_set_multiline(ph, [label, desc], header_pt=FONT_SECTION, font_pt=FONT_BODY)
    for name in ("TextBox 23", "TextBox 29"):
        shape = find_shape(slide, name)
        if shape:
            force_set_text(shape, "")


def apply_slide_18(slide, content: dict) -> None:
    set_title_text(find_title(slide), content["title"])
    intro = find_shape(slide, "TextBox 5")
    body = find_shape(slide, "TextBox 4")
    if intro:
        set_single_text(intro, content["intro"], font_pt=FONT_BODY, auto_fit=True)
    if body:
        set_single_text(body, content["body"], font_pt=FONT_BODY, auto_fit=True)
    step_boxes = sorted(
        [s for s in iter_shapes(slide.shapes) if s.has_text_frame and s.name in ("TextBox 10", "TextBox 11", "TextBox 12", "TextBox 13", "TextBox 14")],
        key=lambda s: s.left,
    )
    for box, step in zip(step_boxes, content["steps"]):
        set_single_text(box, step, font_pt=FONT_BODY, auto_fit=True)


def apply_slide_22(slide, content: dict) -> None:
    set_title_text(find_title(slide), content["title"])
    all_ph = [
        s
        for s in iter_shapes(slide.shapes)
        if s.has_text_frame and s.name.startswith("Text Placeholder")
    ]
    title_ph = sorted([s for s in all_ph if inches(s)[3] > 2.5], key=lambda s: (s.top, s.left))
    body_ph = sorted(
        [
            s
            for s in all_ph
            if inches(s)[3] <= 1.5
            and round(s.top / 914400, 2) not in (1.71, 4.89)
        ],
        key=lambda s: (s.top, s.left),
    )
    for ph, (label, _) in zip(title_ph[:4], content["boxes"]):
        force_set_text(ph, label, font_pt=FONT_STAT_TITLE)
    for ph, (_, desc) in zip(body_ph[:4], content["boxes"]):
        force_set_text(ph, truncate(desc, 90), font_pt=FONT_BODY)
    # Remove legacy duplicate acquisition-framework row (lower template copies)
    for ph in all_ph:
        top = round(ph.top / 914400, 2)
        if top in (1.71, 4.89):
            force_set_text(ph, "")
    for name in ("TextBox 14", "Rectangle: Rounded Corners 61", "TextBox 17"):
        shape = find_shape(slide, name)
        if shape:
            force_set_text(shape, "")


def apply_slide_15(slide, content: dict) -> None:
    set_title_text(find_title(slide), content["title"])
    for shape in iter_shapes(slide.shapes):
        if shape.is_placeholder and "SUBTITLE" in str(shape.placeholder_format.type):
            set_multiline_text(shape, content["subtitle"].split("\n"), font_pt=11.0, bullet_from=-1, auto_fit=True)


APPLY_HANDLERS = {
    1: apply_slide_1,
    3: apply_slide_3,
    5: apply_slide_5,
    6: apply_slide_6,
    7: apply_slide_7,
    8: apply_slide_8,
    9: apply_slide_9,
    10: apply_slide_10,
    12: apply_slide_12,
    15: apply_slide_15,
    16: apply_slide_16,
    17: apply_slide_17,
    18: apply_slide_18,
    22: apply_slide_22,
}


def apply_slide_content(slide, slide_num: int, content: dict) -> None:
    handler = APPLY_HANDLERS.get(slide_num)
    if handler:
        handler(slide, content)
    replace_footer(slide)


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    r_id = slide_id_list[index].rId
    prs.part.drop_rel(r_id)
    del slide_id_list[index]


def reorder_slides(prs: Presentation, order_indices: list[int]) -> None:
    slide_id_list = prs.slides._sldIdLst
    elements = [slide_id_list[i] for i in range(len(slide_id_list))]
    for el in list(slide_id_list):
        slide_id_list.remove(el)
    for idx in order_indices:
        slide_id_list.append(elements[idx])


def generate() -> Path:
    summary = load_summary()
    doc = parse_docx(DOCX_PATH)
    content_map = build_content(doc)

    shutil.copy2(TEMPLATE_PPTX, OUTPUT_FALLBACK)
    prs = Presentation(str(OUTPUT_FALLBACK))

    keep_slides = sorted(content_map.keys())
    total = len(prs.slides)
    for slide_num in sorted((i for i in range(1, total + 1) if i not in keep_slides), reverse=True):
        delete_slide(prs, slide_num - 1)

    current_order = sorted(content_map.keys())
    final_order = [1, 3, 5, 12, 16, 9, 10, 6, 7, 8, 17, 18, 22, 15]
    reorder_slides(prs, [current_order.index(n) for n in final_order])

    for display_num, template_slide_num in enumerate(final_order, start=1):
        slide = prs.slides[display_num - 1]
        apply_slide_content(slide, template_slide_num, content_map[template_slide_num])
        replace_slide_number(slide, display_num)

    prs.save(str(OUTPUT_FALLBACK))

    try:
        shutil.copy2(OUTPUT_FALLBACK, OUTPUT_PPTX)
        output_path = OUTPUT_PPTX
    except PermissionError:
        output_path = OUTPUT_FALLBACK

    manifest = {
        "output": str(output_path),
        "source_docx": str(DOCX_PATH),
        "template_id": summary["template_id"],
        "slide_count": len(final_order),
        "slides": [
            {"number": i + 1, "template_slide": n, "title": content_map[n].get("title", "")[:120]}
            for i, n in enumerate(final_order)
        ],
    }
    manifest_path = output_path.with_suffix(".manifest.json")
    manifest_path.write_text(
        json.dumps(manifest, indent=2, ensure_ascii=False), encoding="utf-8"
    )
    if output_path != OUTPUT_FALLBACK:
        OUTPUT_FALLBACK.with_suffix(".manifest.json").write_text(
            json.dumps(manifest, indent=2, ensure_ascii=False), encoding="utf-8"
        )
    return output_path


if __name__ == "__main__":
    print(f"Created: {generate()}")
