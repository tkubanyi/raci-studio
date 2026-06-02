"""Extract PwC brand graphic elements from a PPTX into a reusable JSON catalog."""

from __future__ import annotations

import hashlib
import json
import re
import sys
import zipfile
import xml.etree.ElementTree as ET
from collections import Counter, defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_INCH = 914400


def emu_to_inches(emu: int | None) -> float | None:
    if emu is None:
        return None
    return round(emu / EMU_PER_INCH, 4)


def emu_to_pt(emu: int | None) -> float | None:
    if emu is None:
        return None
    return round(emu / 12700, 2)


def color_to_hex(color) -> str | None:
    if color is None:
        return None
    try:
        rgb = color.rgb
        if rgb is not None:
            return f"#{rgb}"
    except Exception:
        pass
    try:
        if hasattr(color, "theme_color") and color.theme_color is not None:
            return f"theme:{color.theme_color}"
    except Exception:
        pass
    return None


def extract_theme_colors(pptx_path: Path) -> dict[str, str]:
    colors: dict[str, str] = {}
    with zipfile.ZipFile(pptx_path) as z:
        theme_files = sorted(n for n in z.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml"))
        if not theme_files:
            return colors
        root = ET.fromstring(z.read(theme_files[0]))
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        clr_scheme = root.find(".//a:clrScheme", ns)
        if clr_scheme is None:
            return colors
        for child in clr_scheme:
            tag = child.tag.split("}")[-1]
            srgb = child.find("a:srgbClr", ns)
            sys_clr = child.find("a:sysClr", ns)
            if srgb is not None and srgb.get("val"):
                colors[tag] = f"#{srgb.get('val')}"
            elif sys_clr is not None and sys_clr.get("lastClr"):
                colors[tag] = f"#{sys_clr.get('lastClr')}"
    return colors


def extract_font_scheme(pptx_path: Path) -> dict:
    fonts: dict[str, str] = {}
    with zipfile.ZipFile(pptx_path) as z:
        theme_files = sorted(n for n in z.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml"))
        if not theme_files:
            return fonts
        root = ET.fromstring(z.read(theme_files[0]))
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        font_scheme = root.find(".//a:fontScheme", ns)
        if font_scheme is None:
            return fonts
        for child in font_scheme:
            tag = child.tag.split("}")[-1]
            latin = child.find(".//a:latin", ns)
            if latin is not None and latin.get("typeface"):
                fonts[tag] = latin.get("typeface")
    return fonts


def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield shape
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def shape_kind(shape) -> str:
    t = shape.shape_type
    if t == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if t == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return "auto_shape"
    if t == MSO_SHAPE_TYPE.LINE:
        return "line"
    if t == MSO_SHAPE_TYPE.FREEFORM:
        return "freeform"
    if t == MSO_SHAPE_TYPE.PLACEHOLDER:
        return "placeholder"
    if t == MSO_SHAPE_TYPE.TEXT_BOX:
        return "text_box"
    if t == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if t == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if t == MSO_SHAPE_TYPE.GROUP:
        return "group"
    return str(t).split("(")[0].strip().lower()


def fill_color_hex(shape) -> str | None:
    try:
        if shape.fill and shape.fill.type == MSO_FILL.SOLID:
            return color_to_hex(shape.fill.fore_color)
    except Exception:
        pass
    return None


def sample_font(shape) -> dict | None:
    if not shape.has_text_frame:
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            f = run.font
            if f.name or f.size:
                return {
                    "name": f.name,
                    "size_pt": emu_to_pt(f.size) if f.size else None,
                    "bold": f.bold,
                    "italic": f.italic,
                    "color": color_to_hex(f.color) if f.color else None,
                }
    return None


def infer_slide_purpose(slide, slide_index: int) -> str:
    layout = slide.slide_layout.name.lower()
    texts = []
    for shape in iter_shapes(slide.shapes):
        if shape.has_text_frame and shape.text.strip():
            texts.append(shape.text.strip()[:80])
    joined = " ".join(texts).lower()

    if "title slide" in layout or slide_index == 1:
        return "cover"
    if "conclusion" in layout:
        return "closing"
    if "section header" in layout:
        return "section_divider"
    if "stat box" in layout or "four stat" in layout or "six stat" in layout:
        return "stat_blocks"
    if "chart" in layout:
        return "charts"
    if "bio" in layout or "cv" in layout or "team" in layout:
        return "people_bios"
    if "mockup" in layout:
        return "mockup"
    if any(k in joined for k in ("icon", "graphic", "element", "shape", "color", "font", "typography")):
        return "brand_elements"
    if shape_kind_count(slide, "picture") >= 5:
        return "icon_gallery"
    return "content_template"


def shape_kind_count(slide, kind: str) -> int:
    return sum(1 for s in iter_shapes(slide.shapes) if shape_kind(s) == kind)


def build_analysis(pptx_path: Path) -> dict:
    prs = Presentation(str(pptx_path))
    theme_colors = extract_theme_colors(pptx_path)
    font_scheme = extract_font_scheme(pptx_path)

    font_families: set[str] = set()
    font_sizes: set[float] = set()
    fill_colors: set[str] = set()
    line_colors: set[str] = set()
    font_colors: set[str] = set()

    auto_shape_types: Counter = Counter()
    shape_type_counts: Counter = Counter()
    layout_usage: Counter = Counter()

    auto_shape_catalog: dict[str, dict] = {}
    line_catalog: list[dict] = []
    icons_on_slides: list[dict] = []
    slide_catalog: list[dict] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        layout_name = slide.slide_layout.name
        layout_usage[layout_name] += 1

        slide_shapes_summary = Counter()
        slide_title = None
        pictures: list[dict] = []

        for shape in iter_shapes(slide.shapes):
            kind = shape_kind(shape)
            slide_shapes_summary[kind] += 1
            shape_type_counts[kind] += 1

            if kind == "placeholder" and slide_title is None:
                try:
                    if "TITLE" in str(shape.placeholder_format.type) and shape.text.strip():
                        slide_title = shape.text.strip()[:120]
                except Exception:
                    pass
            if slide_title is None and shape.name and shape.name.startswith("Title") and shape.has_text_frame:
                t = shape.text.strip()
                if t:
                    slide_title = t[:120]

            fc = fill_color_hex(shape)
            if fc:
                fill_colors.add(fc)

            try:
                if shape.line and shape.line.fill:
                    lc = color_to_hex(shape.line.fill.fore_color)
                    if lc:
                        line_colors.add(lc)
            except Exception:
                pass

            fs = sample_font(shape)
            if fs:
                if fs.get("name"):
                    font_families.add(fs["name"])
                if fs.get("size_pt"):
                    font_sizes.add(fs["size_pt"])
                if fs.get("color"):
                    font_colors.add(fs["color"])

            if kind == "auto_shape":
                try:
                    ast = str(shape.auto_shape_type)
                except Exception:
                    ast = "unknown"
                auto_shape_types[ast] += 1
                key = ast
                if key not in auto_shape_catalog:
                    auto_shape_catalog[key] = {
                        "auto_shape_type": ast,
                        "count": 0,
                        "example_names": [],
                        "typical_fill_colors": [],
                        "typical_sizes_in": [],
                    }
                entry = auto_shape_catalog[key]
                entry["count"] += 1
                if len(entry["example_names"]) < 5 and shape.name:
                    entry["example_names"].append(shape.name)
                if fc and fc not in entry["typical_fill_colors"] and len(entry["typical_fill_colors"]) < 8:
                    entry["typical_fill_colors"].append(fc)
                sz = (emu_to_inches(shape.width), emu_to_inches(shape.height))
                if len(entry["typical_sizes_in"]) < 6 and sz not in entry["typical_sizes_in"]:
                    entry["typical_sizes_in"].append(list(sz))

            if kind == "line":
                if len(line_catalog) < 40:
                    line_catalog.append({
                        "name": shape.name,
                        "position": {
                            "left_in": emu_to_inches(shape.left),
                            "top_in": emu_to_inches(shape.top),
                            "width_in": emu_to_inches(shape.width),
                            "height_in": emu_to_inches(shape.height),
                        },
                    })

            if kind == "picture":
                pic: dict = {
                    "slide_index": slide_index,
                    "name": shape.name,
                    "position": {
                        "left_in": emu_to_inches(shape.left),
                        "top_in": emu_to_inches(shape.top),
                        "width_in": emu_to_inches(shape.width),
                        "height_in": emu_to_inches(shape.height),
                    },
                }
                try:
                    pic["content_type"] = shape.image.content_type
                    pic["filename"] = shape.image.filename
                except Exception:
                    pass
                pictures.append(pic)
                icons_on_slides.append(pic)

        purpose = infer_slide_purpose(slide, slide_index)
        slide_catalog.append({
            "index": slide_index,
            "layout_name": layout_name,
            "purpose": purpose,
            "title": slide_title,
            "shape_summary": dict(slide_shapes_summary),
            "picture_count": slide_shapes_summary.get("picture", 0),
            "auto_shape_count": slide_shapes_summary.get("auto_shape", 0),
        })

    media_catalog = extract_media_catalog(pptx_path)

    brand_palette = build_brand_palette(theme_colors, fill_colors, line_colors, font_colors)

    layout_definitions = extract_layout_definitions(prs)

    return {
        "template_id": "pwc_graphic_elements_correct",
        "source_file": pptx_path.name,
        "source_path": str(pptx_path),
        "version": "1.0",
        "purpose": "Brand package catalog: fonts, shapes, colors, slide organization, icons",
        "canvas": {
            "width_in": emu_to_inches(prs.slide_width),
            "height_in": emu_to_inches(prs.slide_height),
            "aspect_ratio": "16:9",
        },
        "theme": {
            "colors": theme_colors,
            "font_scheme": font_scheme,
        },
        "typography": {
            "theme_fonts": font_scheme,
            "families_used": sorted(font_families),
            "sizes_pt": sorted(font_sizes),
        },
        "colors": {
            "brand_palette": brand_palette,
            "theme_colors": theme_colors,
            "fill_colors": sorted(c for c in fill_colors if c and not str(c).startswith("theme:")),
            "line_colors": sorted(c for c in line_colors if c and not str(c).startswith("theme:")),
            "font_colors": sorted(c for c in font_colors if c and not str(c).startswith("theme:")),
            "theme_references": sorted(
                c for c in fill_colors | line_colors | font_colors if c and str(c).startswith("theme:")
            ),
        },
        "shapes": {
            "type_counts": dict(shape_type_counts.most_common()),
            "auto_shape_types": dict(auto_shape_types.most_common()),
            "auto_shape_catalog": list(auto_shape_catalog.values()),
            "line_samples": line_catalog[:30],
            "total_auto_shapes": sum(auto_shape_types.values()),
            "total_groups": shape_type_counts.get("group", 0),
            "total_freeforms": shape_type_counts.get("freeform", 0),
        },
        "slide_organization": {
            "slide_count": len(slide_catalog),
            "layouts_by_frequency": dict(layout_usage.most_common()),
            "layout_definitions": layout_definitions,
            "purposes_by_frequency": dict(Counter(s["purpose"] for s in slide_catalog).most_common()),
            "slide_catalog": slide_catalog,
        },
        "icons": {
            "embedded_picture_count": len(icons_on_slides),
            "media_file_count": len(media_catalog),
            "media_by_extension": dict(Counter(m["extension"] for m in media_catalog).most_common()),
            "media_catalog": media_catalog,
            "pictures_by_slide": summarize_pictures_by_slide(icons_on_slides),
            "icon_gallery_slides": [
                s["index"] for s in slide_catalog if s["purpose"] in ("icon_gallery", "brand_elements")
            ],
        },
        "replication_notes": {
            "fonts": "Use theme font_scheme for major/minor; body text typically Arial at 10–16pt from families_used.",
            "colors": "Primary brand orange is theme accent1 (#FD5108 in related decks); use brand_palette tokens.",
            "shapes": "Reference auto_shape_catalog for standard PwC diagram primitives; preserve position boxes from slide_catalog.",
            "icons": "Use media_catalog filenames when extracting from package; match picture positions from pictures_by_slide.",
            "slides": "Pick layout_name from layouts_by_frequency matching purpose (section_divider, stat_blocks, icon_gallery, etc.).",
        },
    }


def build_brand_palette(
    theme: dict[str, str],
    fills: set[str],
    lines: set[str],
    fonts: set[str],
) -> dict:
    return {
        "primary_orange": theme.get("accent1"),
        "orange_mid": theme.get("accent2"),
        "orange_light": theme.get("accent3"),
        "neutral_dark": theme.get("dk1"),
        "neutral_white": theme.get("lt1"),
        "neutral_gray_light": theme.get("lt2"),
        "accent_gray_1": theme.get("accent4"),
        "accent_gray_2": theme.get("accent5"),
        "accent_gray_3": theme.get("accent6"),
        "hyperlink": theme.get("hlink"),
        "common_fills": sorted(c for c in fills if c and c.startswith("#"))[:25],
        "common_font_colors": sorted(c for c in fonts if c and c.startswith("#"))[:15],
    }


def extract_media_catalog(pptx_path: Path) -> list[dict]:
    catalog = []
    with zipfile.ZipFile(pptx_path) as z:
        for name in sorted(z.namelist()):
            if not name.startswith("ppt/media/"):
                continue
            data = z.read(name)
            filename = Path(name).name
            ext = Path(name).suffix.lower().lstrip(".")
            content_type = {
                "png": "image/png",
                "jpg": "image/jpeg",
                "jpeg": "image/jpeg",
                "gif": "image/gif",
                "svg": "image/svg+xml",
                "emf": "image/x-emf",
                "wmf": "image/x-wmf",
            }.get(ext, f"application/{ext}")
            catalog.append({
                "path": name,
                "filename": filename,
                "extension": ext,
                "content_type": content_type,
                "size_bytes": len(data),
                "sha256_prefix": hashlib.sha256(data).hexdigest()[:12],
            })
    return catalog


def summarize_pictures_by_slide(pictures: list[dict]) -> dict[str, list[dict]]:
    by_slide: dict[str, list[dict]] = defaultdict(list)
    for pic in pictures:
        key = str(pic["slide_index"])
        if len(by_slide[key]) < 20:
            by_slide[key].append({
                "name": pic.get("name"),
                "filename": pic.get("filename"),
                "position": pic.get("position"),
                "content_type": pic.get("content_type"),
            })
    return dict(by_slide)


def extract_layout_definitions(prs: Presentation) -> list[dict]:
    seen: set[str] = set()
    layouts = []
    for slide in prs.slides:
        layout = slide.slide_layout
        if layout.name in seen:
            continue
        seen.add(layout.name)
        placeholders = []
        try:
            for ph in layout.placeholders:
                placeholders.append({
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "name": ph.name,
                })
        except Exception:
            pass
        layouts.append({
            "name": layout.name,
            "placeholder_count": len(placeholders),
            "placeholders": placeholders,
        })
    return sorted(layouts, key=lambda x: x["name"])


def main() -> int:
    if len(sys.argv) < 2:
        pptx = Path(r"C:\Users\tkubanyi001\OneDrive - PwC\Documents\Brand package\Graphic elements_correct.pptx")
        out = Path(__file__).resolve().parent.parent / "data" / "presentation-templates" / "graphic_elements_correct.brand.json"
    else:
        pptx = Path(sys.argv[1])
        out = Path(sys.argv[2]) if len(sys.argv) > 2 else pptx.with_suffix(".brand.json")

    if not pptx.exists():
        print(f"File not found: {pptx}")
        return 1

    analysis = build_analysis(pptx)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(analysis, indent=2, ensure_ascii=False), encoding="utf-8")

    summary_path = out.with_suffix(".summary.json")
    summary = {
        k: analysis[k]
        for k in (
            "template_id",
            "source_file",
            "source_path",
            "canvas",
            "theme",
            "typography",
            "colors",
            "shapes",
            "slide_organization",
            "icons",
            "replication_notes",
        )
        if k in analysis
    }
    summary["slide_organization"] = {
        "slide_count": analysis["slide_organization"]["slide_count"],
        "layouts_by_frequency": analysis["slide_organization"]["layouts_by_frequency"],
        "purposes_by_frequency": analysis["slide_organization"]["purposes_by_frequency"],
        "layout_names": [l["name"] for l in analysis["slide_organization"]["layout_definitions"]],
    }
    summary["icons"] = {
        "embedded_picture_count": analysis["icons"]["embedded_picture_count"],
        "media_file_count": analysis["icons"]["media_file_count"],
        "media_by_extension": analysis["icons"]["media_by_extension"],
        "icon_gallery_slides": analysis["icons"]["icon_gallery_slides"],
    }
    summary_path.write_text(json.dumps(summary, indent=2, ensure_ascii=False), encoding="utf-8")

    print(f"Slides: {analysis['slide_organization']['slide_count']}")
    print(f"Media/icons: {analysis['icons']['media_file_count']}")
    print(f"Auto shapes: {analysis['shapes']['total_auto_shapes']}")
    print(f"Saved: {out}")
    print(f"Summary: {summary_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
