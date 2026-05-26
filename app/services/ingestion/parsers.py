"""Format-specific text and image extraction."""

from __future__ import annotations

import re
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

import pdfplumber
from docx import Document as DocxDocument
from PIL import Image


def extract_txt_md(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def extract_docx(path: Path) -> str:
    doc = DocxDocument(str(path))
    parts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_pdf(path: Path, image_dir: Path) -> tuple[str, list[Path]]:
    parts: list[str] = []
    images: list[Path] = []
    try:
        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(text.strip())
                try:
                    for j, img in enumerate(page.images or []):
                        # pdfplumber image bbox only — use pymupdf fallback below
                        pass
                except Exception:
                    pass
    except Exception:
        parts = []

    text = "\n\n".join(parts)
    if len(text.strip()) < 80:
        images.extend(_pdf_pages_as_images(path, image_dir))
        if not text.strip():
            text = "[PDF with limited text — visual pages sent to AI/OCR pipeline]"

    return text, images


def _pdf_pages_as_images(path: Path, image_dir: Path) -> list[Path]:
    images: list[Path] = []
    try:
        import fitz  # pymupdf

        image_dir.mkdir(parents=True, exist_ok=True)
        doc = fitz.open(str(path))
        for i in range(min(len(doc), 5)):
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            out = image_dir / f"{path.stem}_page{i + 1}.png"
            pix.save(str(out))
            images.append(out)
        doc.close()
    except Exception:
        pass
    return images


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide_no, slide in enumerate(prs.slides, 1):
        slide_bits: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_bits.append(shape.text.strip())
        if slide_bits:
            parts.append(f"Slide {slide_no}: " + " | ".join(slide_bits))
    return "\n".join(parts)


def extract_xlsx(path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    rows: list[str] = []
    for sheet in wb.worksheets:
        rows.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
    return "\n".join(rows)


def extract_vsdx(path: Path) -> str:
    """Extract labels and connector hints from Visio (.vsdx) OOXML package."""
    parts: list[str] = []
    ns_skip = re.compile(r"\{.*\}")
    with zipfile.ZipFile(path) as zf:
        page_names = sorted(n for n in zf.namelist() if n.startswith("visio/pages/") and n.endswith(".xml"))
        for pname in page_names[:20]:
            try:
                root = ET.fromstring(zf.read(pname))
            except ET.ParseError:
                continue
            for elem in root.iter():
                tag = ns_skip.sub("", elem.tag)
                if tag in ("Text", "cp", "pp") and elem.text and elem.text.strip():
                    t = elem.text.strip()
                    if len(t) > 1 and t not in parts[-3:]:
                        parts.append(t)
        for name in zf.namelist():
            if "masters" in name and name.endswith(".xml"):
                try:
                    root = ET.fromstring(zf.read(name))
                    for elem in root.iter():
                        if elem.text and elem.text.strip() and len(elem.text.strip()) > 2:
                            parts.append(elem.text.strip())
                except ET.ParseError:
                    continue
    unique = list(dict.fromkeys(parts))
    header = f"[Visio diagram: {path.name} — {len(unique)} shape labels extracted]\n"
    return header + "\n".join(unique[:500])


def extract_image(path: Path, image_dir: Path) -> tuple[str, list[Path]]:
    image_dir.mkdir(parents=True, exist_ok=True)
    copy = image_dir / path.name
    if copy.resolve() != path.resolve():
        Image.open(path).save(copy)
    text = ""
    try:
        import pytesseract

        text = pytesseract.image_to_string(Image.open(path)) or ""
    except Exception:
        text = "[Image — enable OPENAI_API_KEY for vision analysis or install Tesseract for local OCR]"
    return text.strip(), [copy if copy.exists() else path]
