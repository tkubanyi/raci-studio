"""Generate up to 15-slide brand-formatted decks from structured document data."""

from __future__ import annotations

import json
import re
import shutil
import tempfile
import zipfile
from dataclasses import dataclass
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from presentation_studio.config import PresentationSettings, get_settings

REPO_ROOT = Path(__file__).resolve().parent.parent

FOOTER_DEFAULT = "Presentation Studio | PwC"


@dataclass
class DeckBuildOptions:
    output_path: Path
    brand_pptx: Path
    brand_json: Path
    footer: str = FOOTER_DEFAULT
    client_line: str = "PwC"


_ACTIVE_FOOTER = FOOTER_DEFAULT
_ACTIVE_CLIENT_LINE = "PwC"
SLIDE_W = 13.3333
MARGIN_L = 0.44
MARGIN_R = 0.44
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R

PAD_H = 0.10
PAD_V = 0.06
LINE_LEADING = 1.22
PARA_GAP = 0.04
MIN_BOX_H = 0.38
ICON_ABOVE_GAP = 0.06
MEASURE_SAFETY = 1.14
CONTENT_BOTTOM = 6.78

FONT_TITLE = 28
FONT_SUBTITLE = 15
FONT_SECTION = 14
FONT_CARD_HEAD = 13
FONT_BODY = 12
FONT_BULLET = 11
FONT_SMALL = 10


def load_brand(brand_json: Path) -> dict:
    return json.loads(brand_json.read_text(encoding="utf-8"))


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def truncate(text: str, max_len: int) -> str:
    text = re.sub(r"\s+", " ", text.strip())
    if len(text) <= max_len:
        return text
    cut = text[: max_len - 1].rsplit(" ", 1)[0]
    return cut + "…"


def char_width_in(font_pt: float, bold: bool = False) -> float:
    factor = 0.54 if bold else 0.48
    return (font_pt / 72.0) * factor


def line_height_in(font_pt: float) -> float:
    return (font_pt / 72.0) * LINE_LEADING + PARA_GAP


def wrap_line_count(text: str, width_in: float, font_pt: float, bold: bool = False) -> int:
    if not text or width_in <= 0:
        return 0
    max_chars = max(6, int(width_in / char_width_in(font_pt, bold)))
    words = text.split()
    if not words:
        return 1
    lines, current = 1, 0
    for word in words:
        w = len(word) + (1 if current else 0)
        if current + w > max_chars:
            lines += 1
            current = len(word)
        else:
            current += w
    return lines


def measure_blocks_height(blocks: list[tuple[str, int, bool]], inner_width_in: float) -> float:
    if not blocks:
        return MIN_BOX_H
    total = PAD_V * 2
    for text, size, bold in blocks:
        total += wrap_line_count(text, inner_width_in, size, bold) * line_height_in(size)
    return max(MIN_BOX_H, total * MEASURE_SAFETY)


def uniform_panel_height(
    card_specs: list[tuple[float, list[tuple[str, int, bool]]]],
    *,
    max_height: float | None = None,
) -> float:
    """Max measured height for a set of (width, blocks) specs; optional cap."""
    heights = [
        measure_blocks_height(blocks, width - 2 * PAD_H) for width, blocks in card_specs
    ]
    h = max(heights) if heights else MIN_BOX_H
    if max_height is not None:
        h = min(h, max_height)
    return h


def fit_blocks_to_inner_height(
    blocks: list[tuple[str, int, bool]],
    inner_width_in: float,
    max_inner_h: float,
) -> list[tuple[str, int, bool]]:
    """Trim body lines so content fits within max_inner_h (keeps headings)."""
    if measure_blocks_height(blocks, inner_width_in) <= max_inner_h:
        return blocks
    trimmed = list(blocks)
    while len(trimmed) > 1 and measure_blocks_height(trimmed, inner_width_in) > max_inner_h:
        trimmed.pop()
    if measure_blocks_height(trimmed, inner_width_in) <= max_inner_h:
        return trimmed
    head, tail = trimmed[0], trimmed[-1]
    text, size, bold = tail
    for n in range(len(text), 20, -8):
        candidate = head, (truncate(text, n), size, bold)
        if measure_blocks_height(list(candidate), inner_width_in) <= max_inner_h:
            return list(candidate)
    return [head, (truncate(text, 40), size, bold)]


def content_top_y(has_subtitle: bool = True) -> float:
    return 1.78 if has_subtitle else 1.58


def icon_size_pt(col_width_in: float) -> float:
    return min(0.52, max(0.38, col_width_in * 0.38))


def place_snug_card(
    slide,
    left: float,
    top: float,
    width: float,
    blocks: list[tuple[str, int, bool]],
    fill: str,
    *,
    panel_height: float | None = None,
) -> float:
    """Text panel; use panel_height for uniform sizing across a slide."""
    inner_w = width - 2 * PAD_H
    natural = measure_blocks_height(blocks, inner_w)
    box_h = panel_height if panel_height is not None else natural
    add_panel(slide, left, top, width, box_h, fill)
    add_textbox(slide, left + PAD_H, top + PAD_V, inner_w, box_h - 2 * PAD_V, blocks)
    return box_h


def place_uniform_icon_row(
    slide,
    top: float,
    columns: list[tuple[float, float, list[tuple[str, int, bool]], str, Path | None]],
) -> float:
    """
    columns: (left, width, blocks, fill, icon_path)
    Same panel height per row; icons in a band above panels.
    """
    if not columns:
        return 0.0
    specs = [(width, blocks) for _, width, blocks, _, _ in columns]
    panel_h = uniform_panel_height(specs)
    max_inner = panel_h - 2 * PAD_V
    icon_band = 0.0
    for _, width, _, _, icon in columns:
        if icon:
            icon_band = max(icon_band, icon_size_pt(width) + ICON_ABOVE_GAP)
    card_top = top + icon_band
    for left, width, blocks, fill, icon in columns:
        if icon:
            iw = icon_size_pt(width)
            icon_y = top + (icon_band - ICON_ABOVE_GAP - iw)
            add_picture(slide, icon, left + (width - iw) / 2, icon_y, iw)
        fitted = fit_blocks_to_inner_height(blocks, width - 2 * PAD_H, max_inner)
        place_snug_card(slide, left, card_top, width, fitted, fill, panel_height=panel_h)
    return icon_band + panel_h


def place_uniform_grid(
    slide,
    top: float,
    cells: list[tuple[float, float, list[tuple[str, int, bool]], str]],
    *,
    cols: int,
    gap_x: float,
    gap_y: float,
    max_panel_height: float | None = None,
) -> None:
    """
    cells: (left, width, blocks, fill) in row-major order.
    All panels on the slide share one height.
    """
    if not cells:
        return
    specs = [(width, blocks) for _, width, blocks, _ in cells]
    n_rows = (len(cells) + cols - 1) // cols
    if max_panel_height is None and n_rows > 0:
        available = CONTENT_BOTTOM - top - 0.15
        max_panel_height = (available - gap_y * (n_rows - 1)) / n_rows

    panel_h = uniform_panel_height(specs, max_height=max_panel_height)
    max_inner = panel_h - 2 * PAD_V

    for idx, (left, width, blocks, fill) in enumerate(cells):
        inner_w = width - 2 * PAD_H
        fitted = fit_blocks_to_inner_height(blocks, inner_w, max_inner)
        row, col = divmod(idx, cols)
        y = top + row * (panel_h + gap_y)
        place_snug_card(slide, left, y, width, fitted, fill, panel_height=panel_h)


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    r_id = slide_id_list[index].rId
    prs.part.drop_rel(r_id)
    del slide_id_list[index]


def layout_by_name(prs: Presentation, name: str) -> int:
    for i, layout in enumerate(prs.slide_layouts):
        if layout.name == name:
            return i
    raise KeyError(f"Layout not found: {name}")


def set_run_font(run, *, name="Arial", size=12, bold=False, color="#000000") -> None:
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def configure_tf(tf, *, anchor=MSO_ANCHOR.TOP, tight: bool = True) -> None:
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE if tight else MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = anchor
    m = 0.03 if tight else 0.05
    tf.margin_left = Inches(m)
    tf.margin_right = Inches(m)
    tf.margin_top = Inches(m)
    tf.margin_bottom = Inches(m)


def write_text_frame(tf, blocks: list[tuple[str, int, bool]]) -> None:
    tf.clear()
    for i, (text, size, bold) in enumerate(blocks):
        if not text:
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.space_after = Pt(2)
        if p.runs:
            set_run_font(p.runs[0], size=size, bold=bold)
        p.level = 1 if text.startswith("•") else 0


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    blocks: list[tuple[str, int, bool]],
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    configure_tf(box.text_frame)
    write_text_frame(box.text_frame, blocks)


def add_panel(slide, left: float, top: float, width: float, height: float, fill: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.fill.background()


def add_accent_bar(slide, colors: dict) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(colors["primary_orange"])
    bar.line.fill.background()


def add_footer(slide, colors: dict, slide_num: int) -> None:
    add_textbox(slide, MARGIN_L, 7.05, CONTENT_W, 0.25, [(_ACTIVE_FOOTER, 9, False)])
    add_textbox(slide, SLIDE_W - 0.55, 7.05, 0.4, 0.25, [(str(slide_num), 9, False)])


def set_slide_title(slide, title: str, subtitle: str | None = None) -> None:
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        ph_type = str(shape.placeholder_format.type)
        if "TITLE" in ph_type or "CENTER_TITLE" in ph_type:
            tf = shape.text_frame
            configure_tf(tf, tight=False)
            blocks = [(title.replace("\n", " — "), FONT_TITLE, True)]
            if subtitle:
                blocks.append((subtitle, FONT_SUBTITLE, False))
            write_text_frame(tf, blocks)
            return
    add_textbox(slide, MARGIN_L, 0.44, CONTENT_W, 1.1, [(title, FONT_TITLE, True)])
    if subtitle:
        add_textbox(slide, MARGIN_L, 1.35, CONTENT_W, 0.5, [(subtitle, FONT_SUBTITLE, False)])


def clear_non_placeholder_shapes(slide) -> None:
    to_remove = [s for s in slide.shapes if not s.is_placeholder]
    for shape in to_remove:
        sp = shape.element
        sp.getparent().remove(sp)


def extract_icon_paths(pptx_path: Path, count: int = 16) -> list[Path]:
    dest = Path(tempfile.mkdtemp(prefix="pwc_icons_"))
    icons: list[tuple[int, Path]] = []
    with zipfile.ZipFile(pptx_path) as z:
        for name in z.namelist():
            if not name.startswith("ppt/media/") or not name.lower().endswith(".png"):
                continue
            data = z.read(name)
            if 800 < len(data) < 120000:
                out = dest / Path(name).name
                out.write_bytes(data)
                icons.append((len(data), out))
    icons.sort(key=lambda x: x[0])
    step = max(1, len(icons) // count)
    picked = [p for _, p in icons[::step]][:count]
    return picked if picked else [p for _, p in icons[:count]]


def add_picture(slide, path: Path, left: float, top: float, size: float) -> None:
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), height=Inches(size))


def build_cover(slide, doc: dict, colors: dict) -> None:
    set_slide_title(
        slide,
        doc["title"],
        f"{doc['subtitle']}\n{_ACTIVE_CLIENT_LINE} | {date.today().strftime('%d %B %Y')}",
    )
    add_accent_bar(slide, colors)


def build_section(slide, title: str, colors: dict) -> None:
    set_slide_title(slide, title, None)
    add_accent_bar(slide, colors)


def build_problem_context(slide, doc: dict, colors: dict) -> None:
    set_slide_title(slide, "What? — Current State", "Fragmented SSC and JV operating model")
    top = content_top_y(True)
    paras = doc["what"].get("paragraphs", [])
    ctx = paras[0] if paras else ""
    intent = paras[1] if len(paras) > 1 else ""
    blocks = [
        ("Operating context", FONT_SECTION, True),
        (truncate(ctx, 420), FONT_BODY, False),
        ("Discovery intent", FONT_SECTION, True),
        (truncate(intent, 380), FONT_BODY, False),
    ]
    place_snug_card(slide, MARGIN_L, top, CONTENT_W, blocks, colors["panel_warm_1"])


def build_problem_objective(slide, doc: dict, colors: dict) -> None:
    set_slide_title(slide, "What? — Objectives & Methods", "Fact-based view and proven methodologies")
    top = content_top_y(True)
    paras = doc["what"].get("paragraphs", [])
    obj = paras[2] if len(paras) > 2 else ""
    methods = paras[3] if len(paras) > 3 else ""
    gap = 0.14
    col_w = (CONTENT_W - gap) / 2
    place_snug_card(
        slide,
        MARGIN_L,
        top,
        col_w,
        [("Key objective", FONT_SECTION, True), (truncate(obj, 360), FONT_BODY, False)],
        colors["panel_light"],
    )
    place_snug_card(
        slide,
        MARGIN_L + col_w + gap,
        top,
        col_w,
        [("Methodologies applied", FONT_SECTION, True), (truncate(methods, 360), FONT_BODY, False)],
        colors["panel_warm_1"],
    )


def build_opportunities_grid(
    slide,
    pairs: list[tuple[str, str]],
    colors: dict,
    icons: list[Path],
    *,
    cols: int,
    start: int,
    title: str = "So What? — Opportunities",
    subtitle: str = "Cost, governance and scalability improvements",
) -> None:
    set_slide_title(slide, title, subtitle)
    top = content_top_y(True)
    items = pairs[start : start + cols]
    if not items:
        return
    gap = 0.14
    col_w = (CONTENT_W - gap * (cols - 1)) / cols
    fills = [colors["panel_warm_1"], colors["panel_light"], colors["panel_warm_2"], colors["panel_light"]]
    columns: list[tuple[float, float, list[tuple[str, int, bool]], str, Path | None]] = []
    for i, (opp_title, body) in enumerate(items):
        left = MARGIN_L + i * (col_w + gap)
        blocks = [
            (truncate(opp_title, 48), FONT_CARD_HEAD, True),
            (truncate(body, 110), FONT_BODY, False),
        ]
        icon = icons[i] if i < len(icons) else None
        columns.append((left, col_w, blocks, fills[i % len(fills)], icon))
    place_uniform_icon_row(slide, top, columns)


def build_workstreams(slide, workstreams: list[tuple[str, str]], colors: dict) -> None:
    set_slide_title(slide, "Now What? — Five Workstreams", "Integrated discovery program structure")
    top = content_top_y(True)
    if not workstreams:
        workstreams = [
            ("Service & Process Mapping", "Document services across SSC, JVs and HQ."),
            ("RACI & Governance Mapping", "Define roles and decision rights."),
            ("P&L, Cost & Headcount Review", "Analyze allocations and workforce."),
            ("Opportunity & Placement Logic", "Evaluate criteria for future location."),
            ("Future-State Blueprint Inputs", "Playbook and deployment template inputs."),
        ]
    n = len(workstreams)
    gap = 0.10
    box_w = (CONTENT_W - gap * (n - 1)) / n
    orange_scale = ["#FD5108", "#FE7C39", "#FFAA72", "#FFCDA8", "#FFE8D4"]
    ws_specs: list[tuple[float, list[tuple[str, int, bool]]]] = []
    ws_layout: list[tuple[float, list[tuple[str, int, bool]], str]] = []
    for i, (title, desc) in enumerate(workstreams):
        left = MARGIN_L + i * (box_w + gap)
        blocks = [
            (f"{i + 1}", 15, True),
            (truncate(title, 36), FONT_CARD_HEAD, True),
            (truncate(desc, 95), FONT_BULLET, False),
        ]
        ws_specs.append((box_w, blocks))
        ws_layout.append((left, blocks, orange_scale[i % 5]))
    panel_h = uniform_panel_height(ws_specs)
    max_inner = panel_h - 2 * PAD_V
    for left, blocks, fill in ws_layout:
        fitted = fit_blocks_to_inner_height(blocks, box_w - 2 * PAD_H, max_inner)
        place_snug_card(slide, left, top, box_w, fitted, fill, panel_height=panel_h)


def build_timeline(slide, phases: list[list[str]], colors: dict) -> None:
    set_slide_title(slide, "Tentative Timeline", "8–16 week diagnostic program")
    top = content_top_y(True)
    if not phases:
        phases = [
            ["Mobilization", "Weeks 1–2", "Scope, governance, templates"],
            ["Current-State Discovery", "Weeks 3–6", "Interviews, mapping, RACI"],
            ["Analysis", "Weeks 7–10", "Duplication, best practice, cost"],
            ["Future-State Options", "Weeks 11–14", "Prioritization and placement"],
            ["Handover", "Weeks 15–16", "Findings and next-phase inputs"],
        ]
    phases = phases[:5]
    pentagon_colors = ["#FD5108", "#FE7C39", "#FFAA72", "#FFCDA8", "#FFE8D4"]
    n = len(phases)
    col_w = CONTENT_W / n
    badge_h = line_height_in(FONT_BULLET) + PAD_V
    card_top = top + badge_h + 0.08
    card_w = col_w - 0.08
    tl_specs: list[tuple[float, list[tuple[str, int, bool]]]] = []
    tl_layout: list[tuple[float, list[tuple[str, int, bool]]]] = []
    for i, phase in enumerate(phases):
        left = MARGIN_L + i * col_w
        week = phase[1].replace("Weeks ", "Wk ").replace("–", "-")
        inner = col_w - 0.16
        badge_w = min(inner, char_width_in(FONT_BULLET, True) * len(week) + 0.32)
        badge_left = left + (col_w - badge_w) / 2
        badge = slide.shapes.add_shape(
            MSO_SHAPE.PENTAGON, Inches(badge_left), Inches(top), Inches(badge_w), Inches(badge_h)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = rgb(pentagon_colors[i % 5])
        badge.line.fill.background()
        configure_tf(badge.text_frame)
        write_text_frame(badge.text_frame, [(week, FONT_BULLET, True)])

        blocks = [
            (truncate(phase[0], 32), FONT_CARD_HEAD, True),
            (truncate(phase[2], 85), FONT_BULLET, False),
        ]
        tl_specs.append((card_w, blocks))
        tl_layout.append((left + 0.04, blocks))
    panel_h = uniform_panel_height(tl_specs)
    max_inner = panel_h - 2 * PAD_V
    for left, blocks in tl_layout:
        fitted = fit_blocks_to_inner_height(blocks, card_w - 2 * PAD_H, max_inner)
        place_snug_card(slide, left, card_top, card_w, fitted, colors["panel_light"], panel_height=panel_h)


def build_team(slide, roles: list[tuple[str, str]], colors: dict) -> None:
    set_slide_title(slide, "Required Project Team", "Cross-functional representation from GP Europe")
    top = content_top_y(True)
    if not roles:
        roles = [
            ("Executive Sponsor (1 FTE)", "Strategic direction and priorities"),
            ("Project Lead / PM (2 FTE)", "Delivery, governance, reporting"),
            ("SSC Operations Lead (0.2 FTE)", "Prague SSC capabilities and capacity"),
            ("JV Process Owners (1.2 FTE)", "Local process and cost inputs"),
            ("Finance / FP&A Lead (0.5 FTE)", "P&L and business case"),
            ("HR / Workforce Lead (0.5 FTE)", "Headcount and workforce implications"),
            ("Process Excellence Lead (1 FTE)", "Mapping and standardization"),
            ("Technology Lead (0.5 FTE)", "Systems, data flows, automation"),
            ("Risk, Compliance & Legal (0.5 FTE)", "Regulatory and control validation"),
        ]
    gap_x, gap_y = 0.12, 0.10
    col_w = (CONTENT_W - 2 * gap_x) / 3
    team_cells: list[tuple[float, float, list[tuple[str, int, bool]], str]] = []
    for i, (role, desc) in enumerate(roles[:9]):
        row, col = divmod(i, 3)
        left = MARGIN_L + col * (col_w + gap_x)
        blocks = [
            (truncate(role, 42), FONT_CARD_HEAD, True),
            (truncate(desc, 75), FONT_BULLET, False),
        ]
        fill = colors["panel_light"] if (row + col) % 2 else colors["panel_warm_1"]
        team_cells.append((left, col_w, blocks, fill))
    place_uniform_grid(slide, top, team_cells, cols=3, gap_x=gap_x, gap_y=gap_y)


def build_outputs(slide, outputs: list[str], colors: dict) -> None:
    set_slide_title(slide, "Expected Outputs", "Deliverables at end of discovery")
    top = content_top_y(True)
    if not outputs:
        outputs = [
            "Current-state service and process inventory",
            "Process maps for priority services",
            "RACI matrix and decision rights",
            "P&L, cost and headcount baseline",
            "Duplication and fragmentation analysis",
            "JV capability assessment for scale-up",
            "Service placement recommendations",
            "Opportunity backlog and playbook inputs",
        ]
    col_gap = 0.14
    col_w = (CONTENT_W - col_gap) / 2
    half = (len(outputs) + 1) // 2

    def bullet_blocks(items: list[str], heading: str) -> list[tuple[str, int, bool]]:
        blocks: list[tuple[str, int, bool]] = [(heading, FONT_SECTION, True)]
        for item in items:
            blocks.append((f"• {truncate(item, 62)}", FONT_BULLET, False))
        return blocks

    left_blocks = bullet_blocks(outputs[:half], "Discovery deliverables")
    right_blocks = bullet_blocks(outputs[half:], "Transition & playbook inputs")
    panel_h = uniform_panel_height([(col_w, left_blocks), (col_w, right_blocks)])
    max_inner = panel_h - 2 * PAD_V
    place_snug_card(
        slide,
        MARGIN_L,
        top,
        col_w,
        fit_blocks_to_inner_height(left_blocks, col_w - 2 * PAD_H, max_inner),
        colors["panel_light"],
        panel_height=panel_h,
    )
    place_snug_card(
        slide,
        MARGIN_L + col_w + col_gap,
        top,
        col_w,
        fit_blocks_to_inner_height(right_blocks, col_w - 2 * PAD_H, max_inner),
        colors["panel_warm_1"],
        panel_height=panel_h,
    )


def build_methodologies(slide, groups: list[tuple[str, list[str]]], colors: dict) -> None:
    set_slide_title(slide, "Methodologies & Techniques", "Eight structured discovery approaches")
    top = content_top_y(True)
    gap_x, gap_y = 0.10, 0.10
    col_w = (CONTENT_W - gap_x) / 2
    cells: list[tuple[float, float, list[tuple[str, int, bool]], str]] = []
    for i, (label, bullets) in enumerate(groups[:8]):
        _, col = divmod(i, 2)
        left = MARGIN_L + col * (col_w + gap_x)
        blocks: list[tuple[str, int, bool]] = [(truncate(label, 44), FONT_CARD_HEAD, True)]
        for b in bullets[:2]:
            blocks.append((f"• {truncate(b, 78)}", FONT_SMALL, False))
        if len(blocks) == 1:
            blocks.append(("Applied during discovery workshops", FONT_SMALL, False))
        fill = colors["panel_light"] if i % 2 else colors["panel_warm_1"]
        cells.append((left, col_w, blocks, fill))
    place_uniform_grid(slide, top, cells, cols=2, gap_x=gap_x, gap_y=gap_y)


def build_placement_and_prioritization(
    slide,
    options: list[tuple[str, str]],
    criteria: list[str],
    prioritization: list[str],
    colors: dict,
    icons: list[Path],
) -> None:
    set_slide_title(slide, "Placement & Prioritization", "Future-state options and assessment criteria")
    top = content_top_y(True)
    if not options:
        options = [
            ("Remain local (JV)", "Regulatory or proximity requirements."),
            ("Transition to SSC", "Standardizable processes for Prague."),
            ("Scale JV capability", "Reuse mature JV services."),
            ("Build central", "New scalable central capability."),
        ]
    gap = 0.10
    col_w = (CONTENT_W - 3 * gap) / 4
    fills = [colors["panel_warm_1"], colors["panel_light"], colors["panel_warm_2"], colors["panel_light"]]
    row1_cols: list[tuple[float, float, list[tuple[str, int, bool]], str, Path | None]] = []
    for i, (label, desc) in enumerate(options[:4]):
        left = MARGIN_L + i * (col_w + gap)
        blocks = [(label, FONT_CARD_HEAD, True), (truncate(desc, 72), FONT_SMALL, False)]
        icon = icons[i] if i < len(icons) else None
        row1_cols.append((left, col_w, blocks, fills[i % 4], icon))
    row1_h = place_uniform_icon_row(slide, top, row1_cols)
    row2 = top + row1_h + 0.14
    half_w = (CONTENT_W - gap) / 2
    left_blocks: list[tuple[str, int, bool]] = [("Fit-gap criteria", FONT_SECTION, True)]
    for c in (criteria or [])[:5]:
        left_blocks.append((f"• {c}", FONT_SMALL, False))
    right_blocks: list[tuple[str, int, bool]] = [("Prioritization matrix", FONT_SECTION, True)]
    for p in (prioritization or [])[:5]:
        right_blocks.append((f"• {p}", FONT_SMALL, False))
    panel_h = uniform_panel_height([(half_w, left_blocks), (half_w, right_blocks)])
    max_inner = panel_h - 2 * PAD_V
    place_snug_card(
        slide,
        MARGIN_L,
        row2,
        half_w,
        fit_blocks_to_inner_height(left_blocks, half_w - 2 * PAD_H, max_inner),
        colors["panel_light"],
        panel_height=panel_h,
    )
    place_snug_card(
        slide,
        MARGIN_L + half_w + gap,
        row2,
        half_w,
        fit_blocks_to_inner_height(right_blocks, half_w - 2 * PAD_H, max_inner),
        colors["panel_warm_1"],
        panel_height=panel_h,
    )


def build_closing(slide, doc: dict, colors: dict) -> None:
    set_slide_title(slide, "Thank you", None)
    msg = truncate(doc.get("closing", ""), 280)
    blocks = [
        (msg, FONT_BODY, False),
        ("© 2026 PricewaterhouseCoopers. All rights reserved.", FONT_SMALL, False),
    ]
    inner_w = CONTENT_W - 2 * PAD_H
    box_h = measure_blocks_height(blocks, inner_w)
    top = 4.85 - box_h / 2
    add_textbox(slide, MARGIN_L, top, CONTENT_W, box_h, blocks)
    add_accent_bar(slide, colors)


def send_shapes_to_back(slide) -> None:
    sp_tree = slide.shapes._spTree
    shapes = list(slide.shapes)
    panels, others = [], []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if not shape.has_text_frame:
                panels.append(shape)
            elif not shape.text_frame.text.strip():
                panels.append(shape)
            else:
                others.append(shape)
        else:
            others.append(shape)
    for shape in panels + others:
        sp_tree.remove(shape.element)
        sp_tree.append(shape.element)


def align_slide(slide, slide_num: int, colors: dict) -> None:
    send_shapes_to_back(slide)
    add_footer(slide, colors, slide_num)


def build_deck(doc: dict, options: DeckBuildOptions) -> Path:
    global _ACTIVE_FOOTER, _ACTIVE_CLIENT_LINE
    _ACTIVE_FOOTER = options.footer
    _ACTIVE_CLIENT_LINE = options.client_line

    brand = load_brand(options.brand_json)
    palette = brand["colors"]["brand_palette"]
    colors = {
        "primary_orange": palette["primary_orange"],
        "panel_light": "#F5F7F8",
        "panel_warm_1": "#FFF5ED",
        "panel_warm_2": "#FFE8D4",
    }

    icons = extract_icon_paths(options.brand_pptx, 20)
    opps = doc["opportunities"]

    options.output_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(options.brand_pptx, options.output_path)
    prs = Presentation(str(options.output_path))
    while len(prs.slides) > 0:
        delete_slide(prs, 0)

    layouts = {
        "title": layout_by_name(prs, "Title Slide"),
        "section": layout_by_name(prs, "Section Header"),
        "content": layout_by_name(prs, "Title Only"),
        "closing": layout_by_name(prs, "Conclusion 1"),
    }

    so_paras = doc.get("so_what", {}).get("paragraphs") or []
    intro_para = so_paras[0] if so_paras else ""

    slide_builders = [
        (layouts["title"], lambda s: build_cover(s, doc, colors)),
        (layouts["section"], lambda s: build_section(s, "1. What? — Problem Statement", colors)),
        (layouts["content"], lambda s: build_problem_context(s, doc, colors)),
        (layouts["content"], lambda s: build_problem_objective(s, doc, colors)),
        (layouts["section"], lambda s: build_section(s, "2. So What? — Opportunities", colors)),
        (
            layouts["content"],
            lambda s: build_opportunities_grid(
                s,
                opps,
                colors,
                icons[0:3],
                cols=3,
                start=0,
            ),
        ),
        (
            layouts["content"],
            lambda s: build_opportunities_grid(
                s,
                opps,
                colors,
                icons[3:7],
                cols=4,
                start=3,
                title="So What? — Opportunities (cont.)",
                subtitle=truncate(intro_para, 95) if intro_para else "Governance, cost and scalability",
            ),
        ),
        (layouts["section"], lambda s: build_section(s, "3. Now What? — Approach & Delivery", colors)),
        (layouts["content"], lambda s: build_workstreams(s, doc["workstreams"], colors)),
        (layouts["content"], lambda s: build_timeline(s, doc["timeline"], colors)),
        (layouts["content"], lambda s: build_team(s, doc["team_roles"], colors)),
        (layouts["content"], lambda s: build_outputs(s, doc["outputs"], colors)),
        (layouts["content"], lambda s: build_methodologies(s, doc["methodology_groups"], colors)),
        (
            layouts["content"],
            lambda s: build_placement_and_prioritization(
                s,
                doc["placement_options"],
                doc["fit_gap_criteria"],
                doc["prioritization"],
                colors,
                icons[8:12],
            ),
        ),
        (layouts["closing"], lambda s: build_closing(s, doc, colors)),
    ]

    for i, (layout_idx, builder) in enumerate(slide_builders, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        clear_non_placeholder_shapes(slide)
        builder(slide)
        align_slide(slide, i, colors)

    prs.save(str(options.output_path))
    return options.output_path


def build_deck_from_settings(
    doc: dict,
    output_path: Path | None = None,
    settings: PresentationSettings | None = None,
) -> Path:
    cfg = settings or get_settings()
    out = output_path or (cfg.default_output_dir / "presentation_output.pptx")
    options = DeckBuildOptions(
        output_path=out,
        brand_pptx=cfg.brand_pptx,
        brand_json=cfg.brand_json,
        footer=cfg.default_footer,
    )
    return build_deck(doc, options)


if __name__ == "__main__":
    from presentation_studio.content_parser import load_docx

    sample = REPO_ROOT / "data" / "presentation-templates"
    docx = sample / "Project_Vienna_Discovery_Phase.pptx"
    print("CLI expects a .docx source; use Presentation Studio app for full workflow.")
